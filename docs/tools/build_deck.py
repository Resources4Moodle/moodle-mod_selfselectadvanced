#!/usr/bin/env python3
"""Build the mod_selfselectadvanced capability deck.

Adapted from mod_teamrecruit's /srv/ci/build_deck.py: same 16:9 layout,
same margins, same five slide-grammar functions, same Liberation Sans
typography (SIL OFL, licence-compatible, installed on the CI box - no
font is embedded or downloaded). Runs on the CI box, where python-pptx
and the screenshots live:

    python3 build_deck.py /srv/ci/selfselectadvanced-howto /srv/ci/out.pptx

The version-footer release and supported Moodle range are read from
the plugin's own version.php next to this script (two directories up),
not hard-coded, so a later release never leaves the footer stale.
"""
import sys
import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

SHOTS = sys.argv[1] if len(sys.argv) > 1 else '/srv/ci/selfselectadvanced-howto'
OUT = sys.argv[2] if len(sys.argv) > 2 else '/srv/ci/mod_selfselectadvanced-capability-deck.pptx'
PLUGINROOT = os.path.normpath(os.path.join(os.path.dirname(os.path.abspath(__file__)), '..', '..'))

INK = RGBColor(0x1B, 0x25, 0x38)
ACCENT = RGBColor(0x0F, 0x6C, 0xBF)
MUTED = RGBColor(0x5A, 0x66, 0x75)
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RULE = RGBColor(0xC8, 0xD1, 0xDA)

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


def branch_str(code):
    """Turn a Moodle $plugin->supported branch code (e.g. 502) into '5.2'."""
    major, minor = divmod(int(code), 100)
    return '%d.%d' % (major, minor)


def read_version():
    """Read release + supported range straight from version.php."""
    path = os.path.join(PLUGINROOT, 'version.php')
    with open(path, encoding='utf-8') as fh:
        text = fh.read()
    release = re.search(r"\$plugin->release\s*=\s*'([^']+)'", text).group(1)
    supported = re.search(r"\$plugin->supported\s*=\s*\[\s*(\d+)\s*,\s*(\d+)\s*\]", text)
    lo, hi = branch_str(supported.group(1)), branch_str(supported.group(2))
    lts = ' LTS' if lo == '4.5' else ''
    return release, 'Moodle %s%s to %s' % (lo, lts, hi)


RELEASE, MOODLERANGE = read_version()


def frame(slide, left, top, width, height):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    return tf


def para(tf, text, size=16, bold=False, color=INK, after=8, first=False, bullet=False):
    p = tf.paragraphs[0] if first else tf.add_paragraph()
    p.text = ('•  ' + text) if bullet else text
    p.space_after = Pt(after)
    for r in p.runs:
        r.font.size = Pt(size)
        r.font.bold = bold
        r.font.color.rgb = color
        # Liberation Sans: SIL OFL, GPL-compatible, and metric-compatible
        # with Arial, so the deck renders sensibly on a machine that has
        # neither. No font is embedded or downloaded.
        r.font.name = 'Liberation Sans'
    return p


def band(slide, height=Inches(1.1), color=ACCENT):
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, height)
    sh.fill.solid()
    sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def heading(slide, title, kicker=None):
    band(slide)
    tf = frame(slide, Inches(0.55), Inches(0.10), Inches(12.2), Inches(0.95))
    para(tf, title, size=26, bold=True, color=WHITE, first=True, after=0)
    if kicker:
        t2 = frame(slide, Inches(0.55), Inches(1.18), Inches(12.2), Inches(0.4))
        para(t2, kicker, size=13, color=MUTED, first=True)


def picture(slide, name, left, top, max_w, max_h):
    path = os.path.join(SHOTS, name)
    if not os.path.exists(path):
        raise SystemExit('missing screenshot: ' + path)
    from PIL import Image
    with Image.open(path) as im:
        w, h = im.size
    scale = min(max_w / w, max_h / h)
    pic = slide.shapes.add_picture(path, left, top, width=Emu(int(w * scale)), height=Emu(int(h * scale)))
    pic.line.color.rgb = RULE
    pic.line.width = Pt(0.75)
    return pic


def title_slide():
    s = prs.slides.add_slide(BLANK)
    band(s, prs.slide_height, ACCENT)
    tf = frame(s, Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.6))
    para(tf, 'Group self-selection (Advanced)', size=48, bold=True, color=WHITE, first=True, after=6)
    para(tf, 'Students form their own project teams, a guide reviews each one, and a team only '
             'becomes official once its composition meets the rules the activity sets.',
         size=21, color=WHITE, after=24)
    para(tf, 'mod_selfselectadvanced %s  ·  %s  ·  capability deck for evaluators' % (RELEASE, MOODLERANGE),
         size=15, color=WHITE)
    return s


def shot_slide(title, kicker, bullets, shot, footnote=None):
    if len(bullets) > 4:
        raise SystemExit('too many bullets on slide: ' + title)
    s = prs.slides.add_slide(BLANK)
    heading(s, title, kicker)
    tf = frame(s, Inches(0.55), Inches(1.75), Inches(4.35), Inches(5.3))
    for i, b in enumerate(bullets):
        para(tf, b, size=14, color=INK, after=10, first=(i == 0), bullet=True)
    if footnote:
        para(tf, footnote, size=11, color=MUTED, after=0)
    picture(s, shot, Inches(5.25), Inches(1.70), Inches(7.6).emu, Inches(5.5).emu)
    return s


def table_slide(title, kicker, headers, rows, widths=None):
    s = prs.slides.add_slide(BLANK)
    heading(s, title, kicker)
    cols = len(headers)
    widths = widths or [1.0 / cols] * cols
    left0, top0 = Inches(0.55), Inches(1.75)
    total_w = Inches(12.2)
    x = left0
    for c, head in enumerate(headers):
        w = Emu(int(total_w * widths[c]))
        tf = frame(s, x, top0, w, Inches(0.4))
        para(tf, head, size=12, bold=True, color=ACCENT, first=True, after=0)
        x = Emu(x + w)
    line = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, left0, Inches(2.18), total_w, Pt(1))
    line.fill.solid(); line.fill.fore_color.rgb = RULE; line.line.fill.background(); line.shadow.inherit = False
    y = Inches(2.32)
    for row in rows:
        x = left0
        tallest = Inches(0.4)
        for c, cell in enumerate(row):
            w = Emu(int(total_w * widths[c]))
            tf = frame(s, x, y, w, Inches(0.4))
            para(tf, cell, size=11.5, color=INK, bold=(c == 0), first=True, after=0)
            est = Inches(0.30) * max(1, int(len(cell) / max(14, int(widths[c] * 95)) + 0.9))
            tallest = max(tallest, est)
            x = Emu(x + w)
        y = Emu(y + tallest)
    return s


def text_slide(title, kicker, pairs, note=None):
    if len(pairs) > 6:
        raise SystemExit('too many items on closing-style slide: ' + title)
    s = prs.slides.add_slide(BLANK)
    heading(s, title, kicker)
    tf = frame(s, Inches(0.7), Inches(1.85), Inches(11.9), Inches(5.2))
    first = True
    for label, body in pairs:
        para(tf, label, size=15, bold=True, color=ACCENT, after=2, first=first)
        para(tf, body, size=13.5, color=INK, after=12)
        first = False
    if note:
        para(tf, note, size=11.5, color=MUTED, after=0)
    return s


def section(number, title, blurb):
    s = prs.slides.add_slide(BLANK)
    band(s, prs.slide_height, LIGHT)
    tf = frame(s, Inches(1.0), Inches(2.7), Inches(11.3), Inches(2.4))
    para(tf, number, size=18, bold=True, color=ACCENT, first=True, after=6)
    para(tf, title, size=38, bold=True, color=INK, after=12)
    para(tf, blurb, size=18, color=MUTED)
    return s


# ---------------------------------------------------------------- deck
title_slide()

text_slide(
    'What this activity does', 'Three sentences, no jargon.',
    [
        ('The idea', 'Students form their own project teams inside Moodle, up to the size the activity '
                     'allows, and a guide reviews each team before it becomes official.'),
        ('Composition, not chance', 'A team only becomes official once it meets the counting rule and seat '
                                    'plan the activity sets, so every team ends up matching the mix a course '
                                    'actually needs, not just whoever happened to ask first.'),
        ('Everything stays inside Moodle', 'No external service is involved at any point - every '
                                           'notification, every review and every list lives in the activity itself.'),
    ],
)

table_slide(
    'Who can do what', 'Every capability is a standard Moodle capability and can be overridden per role. '
    'The guide is this plugin’s own non-editing-teacher role, distinct from the editing teacher who '
    'configures the activity.',
    ['', 'Admin', 'Editing teacher', 'Teacher (guide)', 'Student'],
    [
        ['Add this activity to a course', 'Yes', 'Yes', 'No', 'No'],
        ['Configure rules, quotas, dates and the seat plan', 'Yes', 'Yes', 'No', 'No'],
        ['Create or edit an override for one guide or group', 'Yes', 'Yes', 'No', 'No'],
        ['Bring in or edit student attributes', 'Yes', 'Only if delegated', 'No', 'No'],
        ['Create a group and become its leader', 'Yes', 'No', 'No', 'Yes'],
        ['Accept or decline an invitation or nomination', 'Yes', 'No', 'No', 'Yes'],
        ['Appear on the guide list; review and approve groups', 'Yes', 'Only if delegated', 'Yes', 'No'],
        ['Freeze a group', 'Yes', 'Only if delegated', 'Yes', 'No'],
        ['Unfreeze a group', 'Yes', 'Yes', 'Only if delegated', 'No'],
        ['See every group, member and report', 'Yes', 'Yes', 'Yes', 'No'],
    ],
    widths=[0.30, 0.16, 0.19, 0.19, 0.16],
)

section('Part one', 'The administrator', 'Site-wide vocabulary and participant data, kept outside any single course.')

shot_slide(
    'The department word list', 'Site administration · a two-level tree',
    [
        'Departments and their sub-departments are set once, site-wide, in Site administration.',
        'Every course using this activity draws its counting rules and seat plans from the same list.',
        'A department can be added, renamed, reordered or removed at any time.',
        'Nothing here belongs to a single course - it is shared across the whole site.',
    ],
    '01-departments.png',
)

shot_slide(
    'Bringing in student attributes', 'Site administration · gender, department, sub-department, '
    'mobile, programme',
    [
        'These are the plugin’s own records, never the Moodle profile, and never used to create an account.',
        'A CSV brings in many students at once; a single row can also be added or corrected by hand.',
        'An unrecognised department or programme in a CSV is added to the site list automatically.',
        'Only a site administrator - or a role a site chooses to delegate this to - can reach this page.',
    ],
    '02-attributes.png',
)

section('Part two', 'The teacher', 'Every meaningful setting, shown once, before the first group forms.')

shot_slide(
    'Team size and who may lead', 'Editing teacher · once per activity',
    [
        'Minimum and maximum team size decide when a team may submit, and how large it may grow.',
        'A cap on how many teams one student may lead, separate from a cap on how many they may join.',
        'Grading rides the activity’s own point scale, the same as any other Moodle activity.',
    ],
    '03-settings-group-size.png',
)

shot_slide(
    'Guides, and the pick-that-team window', 'Editing teacher · guide capacity and the listing flow',
    [
        'A guide either volunteers a capacity, or the activity sets one fixed ceiling for everyone.',
        'A response window gives the guide a set time to review, with an optional automatic approval past it.',
        'Turning on team listing lets a forming team be found by guides before it submits; a guide '
        'expresses interest, and the leader always decides.',
        'How many open interests a guide may hold, and whether they see who else is interested, are set here.',
    ],
    '04-settings-guides-eoi.png',
)

shot_slide(
    'Formation window and penalties', 'Editing teacher · dates and the sequence-of-joining grade',
    [
        'An opening date, a decide-by date and a hard cutoff are all optional.',
        'A per-day late penalty applies once a member joins after the decide-by date.',
        'A separate penalty applies to a student who never reaches a team at all.',
        'Every student’s grade shows the step-by-step breakdown behind it.',
    ],
    '05-settings-dates-penalties.png',
)

shot_slide(
    'The counting rule and the seat plan', 'Editing teacher · quotas and seats',
    [
        'A counting rule states a minimum or maximum of one attribute value - at least one Female '
        'member, for instance.',
        'A seat plan books specific seats in order - at least one from Computer Science, in this example.',
        'Every rule and every seat is checked live on the team’s own page, worded as plainly as it is enforced.',
        'A team cannot submit, be approved, or be frozen while either is unmet.',
    ],
    '06-counting-rule-and-seats.png',
)

section('Part three', 'The student forms a team', 'One action per screen, in the order a student meets them.')

shot_slide(
    'Create a group', 'Student · the leader-to-be',
    [
        'Naming the team and describing the project is the whole first step.',
        'The student who creates it is its leader from that moment.',
        'A cap on how many teams a student may lead or join is enforced here, not left to be discovered later.',
    ],
    '07-create-a-group.png',
)

shot_slide(
    'Invite a teammate', 'Student · the leader',
    [
        'An invitation reserves a seat immediately, so a team can never be over-invited.',
        'The invitee sees it waiting the next time they sign in, and can accept or decline.',
        'An unanswered invitation can be withdrawn by the leader at any time.',
    ],
    '08-invite-a-teammate.png',
)

shot_slide(
    'The composition panel', 'Student · every member sees it',
    [
        'Every counting rule and every seat appears here, satisfied or not, in plain language.',
        'A shortfall says exactly what is missing - never just a red mark.',
        'The panel updates the moment membership changes; nobody has to ask a teacher whether the team is ready.',
    ],
    '09-composition-panel.png',
)

shot_slide(
    'List the team for guides', 'Student · the leader, while still forming',
    [
        'Listing is entirely optional, and can be withdrawn at any time before submission.',
        'A listed, still-forming team becomes visible to every guide on the pick-that-team page.',
        'Listing changes nothing about how the team is judged - it only changes who can find it.',
    ],
    '10-list-the-team.png',
)

shot_slide(
    'Interest from guides', 'Student · the leader sees every open interest',
    [
        'A guide’s interest arrives with a short note, never a binding offer.',
        'Every member sees how many guides are interested; only the leader sees who, and decides.',
        'The activity can limit a leader to seeing one open interest at a time, oldest first.',
    ],
    '11-interest-from-guides.png',
)

shot_slide(
    'Accept an interest', 'Student · the leader',
    [
        'Accepting pre-assigns that guide to the team - before the team has even submitted.',
        'Every other open interest in the same team is declined automatically, and each of those guides is told.',
        'The team keeps forming as normal; submission now goes straight to the guide who was accepted.',
    ],
    '12-accept-an-interest.png',
)

shot_slide(
    'Submit for guide review', 'Student · the leader',
    [
        'Submission locks membership - only a manager’s staged move can change it from here.',
        'A required written proposal is checked for at the moment of submission, not after.',
        'The team must already meet its counting rule and seat plan before it can submit at all.',
    ],
    '13-submit-for-review.png',
)

section('Part four', 'The guide', 'Browsing, reviewing, and the record every decision leaves.')

shot_slide(
    'Browse listed teams', 'Guide · pick that team',
    [
        'Every forming team that has chosen to be listed appears here as a card.',
        'A count of open interest sorts the most sought-after teams to the top.',
        'Picking a team sends the leader a note and waits for their decision - it commits nothing on its own.',
    ],
    '14-browse-listed-teams.png',
)

shot_slide(
    'The guide’s own dashboard', 'Guide · a working guide',
    [
        'Four figures: teams currently guided, interests still open, and the guide’s own expired and '
        'declined history.',
        'Each figure links straight to the list behind it.',
        'A guide who has volunteered a capacity sees it here, and can change it at any time.',
    ],
    '15-guide-dashboard.png',
)

shot_slide(
    'The interest list, and the team’s contacts', 'Guide · one accepted interest, drilled in',
    [
        'Every interest a guide has raised, with its outcome, in one sortable list.',
        'Opening an accepted interest lists the team’s members, each with the contact channel they volunteered.',
        'A WhatsApp link appears only for a member who gave a number; one click reaches the whole team at once.',
    ],
    '16-interest-list-and-contacts.png',
)

shot_slide(
    'Review a submission', 'Guide · the assigned reviewer',
    [
        'The written proposal sits inline, not one click away.',
        'Approve is irreversible and asks for confirmation; return sends the team back with a mandatory comment.',
        'The guide’s own working notes are private, and saved separately from the decision.',
    ],
    '17-guide-review.png',
)

section('Part five', 'Firm, and frozen', 'What a team - and its members - see once a decision is made.')

shot_slide(
    'A firm team, with its mark', 'Guide · once approved',
    [
        'Approval is the point a guide’s mark, if one is given, starts counting toward every member’s grade.',
        'A firm team keeps its shape until a manager’s staged move changes it.',
        'The grade breakdown always shows which step of a student’s own sequence this mark belongs to.',
    ],
    '18-firm-with-a-mark.png',
)

shot_slide(
    'A frozen team', 'Student · once frozen',
    [
        'Freezing mirrors the team into a real Moodle course group, so it can be used anywhere groups are.',
        'Only a group this plugin created is ever touched - nothing else in the course.',
        'Unfreezing restores the exact roster this plugin last authorised, discarding only changes made outside it.',
    ],
    '19-frozen-team.png',
)

section('Part six', 'Reports', 'What a teacher reads, and can act on.')

shot_slide(
    'Students with no team yet', 'Editing teacher · the flagged report',
    [
        'Every enrolled student who can respond, cross-checked against every confirmed membership - not '
        'just one team at a time.',
        'A student can be missing from every other list and still show up here.',
        'The same report also flags missing attributes, guide load and quota shortfalls, each on its own tab.',
    ],
    '20-students-with-no-team.png',
)

shot_slide(
    'Formation analytics', 'Editing teacher · funnel and timing',
    [
        'How many teams were created, submitted, made firm and frozen, at a glance.',
        'Median time from creation to submission, and from listing to first interest, so one term can be '
        'compared with the next.',
        'Every figure exports in the site’s own choice of format.',
    ],
    '21-formation-analytics.png',
)

text_slide(
    'Privacy', 'Everything declared through Moodle’s Privacy API.',
    [
        ('What is stored', 'Group membership, invitations, guide decisions, expressions of interest, and '
                           'the participant attributes a site chose to bring in - gender, department, '
                           'sub-department, mobile number, programme.'),
        ('Contact channels are volunteered', 'A mobile number is read from the site’s own ingested '
                                             'attributes, never asked of the student directly by this '
                                             'activity, and a WhatsApp link appears only where a number is on file.'),
        ('What is exported', 'Every export a person triggers themselves - their own team, their own '
                             'history - carries only what that person is already allowed to see.'),
        ('What erasure removes', 'A deleted user’s memberships, invitations, interests and attribute '
                                 'record are all removed; history a team’s other members still need - who '
                                 'joined, when - is anonymised rather than deleted outright.'),
    ],
)

text_slide(
    'Limits', 'What this activity deliberately refuses to do.',
    [
        ('Team size', 'A hard minimum and maximum; nothing in between is negotiable per team.'),
        ('One leader, one cap on memberships', 'A student leads or joins only as many teams as the '
                                               'activity allows - never an unlimited number.'),
        ('The counting rule and seat plan bind every stage', 'Submission, approval and freezing all '
                                                             're-check compliance live; a team cannot slip '
                                                             'through by moving fast.'),
        ('Refused: a guide claiming a team unilaterally', 'Picking a team only ever sends a note - the '
                                                          'leader’s acceptance is what assigns a guide, '
                                                          'never the guide’s own choice alone.'),
        ('Refused: any outbound call', 'No external service, of any kind, at any point - not for '
                                       'suggestions, not for messaging, not for storage.'),
    ],
)

prs.save(OUT)
print('wrote', OUT, os.path.getsize(OUT), 'bytes')
print('release', RELEASE, '|', MOODLERANGE)
